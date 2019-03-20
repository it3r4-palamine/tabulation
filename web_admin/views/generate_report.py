# from ..forms.transaction_types import *
from ..models.transaction_types import *
# from ..forms.company import *
from ..models.assessments import *
from ..forms.assessments import *
from ..models.multiple_choice import *
# from ..forms.generate_report import *
from ..models.company_assessment import *
from ..views.common import *
from django.db.models import *
import os, shutil, errno
from django.conf import settings
import sys, traceback, os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from difflib import SequenceMatcher

def generate_report(request,generate_report_id):
	company_assessment = Company_assessment.objects.get(id=generate_report_id)
	pagename = "Reference No.: " + company_assessment.reference_no + " - " + company_assessment.company.name
	datus = {
		"id" : company_assessment.pk,
		"status" : str(company_assessment.is_complete),
		# "date_from" : company_assessment.date_from,
		# "date_to" : company_assessment.date_to,
		"reference_no" : str(company_assessment.reference_no),
		# "company" : company_assessment.company.get_dict(),
		"transaction_type" : company_assessment.transaction_type
	}
	return render(request, 'generate_report/generate_report.html', {"pagename":pagename,"datus":datus})

def download_dialog(request):
	return render(request, 'generate_report/dialogs/download_dialog.html')

def read_assessments(request):
	try:
		data = req_data(request)
		results = {'company_assessment':[],'data':[],'scores':[]}

		transactionTypeArr = []
		if 'date_from' in data and 'date_to' in data:
			filters = {}
			filters['date__range'] = data['date_from'], data['date_to']
			filters['company_assessment'] = data['id']
			filters['is_deleted'] = False

			assessment_sessions = Assessment_session.objects.filter(**filters)
			for sessions in assessment_sessions:
				if sessions.transaction_type and sessions.transaction_type.pk not in transactionTypeArr:
					transactionTypeArr.append(sessions.transaction_type.pk)
			transaction_types = Exercise.objects.filter(id__in=transactionTypeArr)
		else:
			transaction_types = Exercise.objects.filter(id__in=data['transaction_type'])
		datus = []
		scoresArr = []

		company_assessment_data = Company_assessment.objects.get(id=data['id'])
		results['company_assessment'] = company_assessment_data.get_dict()

		transactionTypesArr = []
		for transaction_type in transaction_types:
			transactionTypesRow = transaction_type.get_dict()
			scores = Assessment_score.objects.filter(transaction_type=transaction_type.pk,company_assessment=data['id'])
			for score in scores:
				score_list = score.get_dict()
				scoresArr.append(score_list)

			questions = Assessment_question.objects.filter(Q(is_active=True,transaction_type=transaction_type.pk) | Q(is_active=True,transaction_types__overlap=[transaction_type.pk])).order_by('transaction_type__name')
			all_findings = []

			for question in questions:
				row = question.get_dict()

				if not question.uploaded_question:
					findings = Assessment_finding.objects.filter(is_active=True,question=question.pk)
					possible_finding = []
					if row['is_general'] == True and row['transaction_types']:
						for types in row['transaction_types']:
							if types['id'] == transaction_type.pk:
								row['transaction_type'] = types

					answers_chosen = None
					try:
						answers_chosen = Assessment_answer.objects.get(company_assessment=data['id'],question=question.pk,is_deleted=False)
					except Assessment_answer.DoesNotExist:
						# continue
						pass
					wrong_answer = Decimal(0)
					chosen = [] 
					if answers_chosen.choice:
						for choices in answers_chosen.choice:
							choices = Choice.objects.get(pk=choices)
							if choices.is_answer == False:
								wrong_answer += 1
							answer_dict = {'id':choices.pk,'value':choices.value,'is_answer':choices.is_answer}
							chosen.append(answer_dict)
						row['answers'] = chosen
					else:
						row['answers'] = answers_chosen.text_answer

					if answers_chosen.document_image:
						row['image'] = str(answers_chosen.document_image)
					else:
						row['image'] = None

					if data['type'] == 'pdf':
						for finding in findings:
							row['dupes'] = False
							if finding.value not in all_findings:
								row['dupes'] = True
								finding_question = finding.get_dict()
								possible_finding.append(finding_question)
								row['findings'] = possible_finding
								
								effects = Assessment_effect.objects.filter(question=question.pk,is_active=True)
								possible_effect = []
								for effect in effects:
									effect_question = effect.get_dict()
									possible_effect.append(effect_question)
								row['effects'] = possible_effect
								
								choices = Choice.objects.filter(question=question.pk,is_active=True)
								answers = []
								for choice in choices:
									answer_choice = choice.get_dict()
									answers.append(answer_choice)
								row['choices'] = answers
						datus.append(row)

						for all_finding in findings:
							if wrong_answer > 0:
								all_findings.append(all_finding.value)
					elif data['type'] == 'ppt':
						for finding in findings:
							# if finding.value not in all_findings:
							finding_question = finding.get_dict()
							possible_finding.append(finding_question)
							row['findings'] = possible_finding
								
						effects = Assessment_effect.objects.filter(question=question.pk,is_active=True)
						possible_effect = []
						for effect in effects:
							effect_question = effect.get_dict()
							possible_effect.append(effect_question)
						row['effects'] = possible_effect
						
						choices = Choice.objects.filter(question=question.pk,is_active=True)
						answers = []
						for choice in choices:
							answer_choice = choice.get_dict()
							answers.append(answer_choice)
						row['choices'] = answers
						datus.append(row)

						for all_finding in findings:
							if wrong_answer > 0:
								all_findings.append(all_finding.value)
				else:
					findings = Assessment_finding.objects.filter(is_active=True,question=question.pk)
					possible_finding = []
					if row['is_general'] == True and row['transaction_types']:
						for types in row['transaction_types']:
							if types['id'] == transaction_type.pk:
								row['transaction_type'] = types

					answers_chosen = None
					try:
						answers_chosen = Assessment_answer.objects.get(company_assessment=data['id'],question=question.pk,is_deleted=False)
					except Assessment_answer.DoesNotExist:
						# continue
						pass
					wrong_answer = Decimal(0)

					image_answers = Assessment_upload_answer.objects.filter(is_deleted=False,question=question.pk,transaction_type=transaction_type.pk,company_assessment=data['id'])
					imageAnswersArr = []

					for image_answer in image_answers:
						for correct_answers in row['answers']:
							if image_answer.item_no == correct_answers['item_no']:
								if image_answer.answer != correct_answers['answer']:
									wrong_answer += 1
					# chosen = [] 
					if answers_chosen:
						if answers_chosen.choice:
							for choices in answers_chosen.choice:
								choices = Choice.objects.get(pk=choices)
								if choices.is_answer == False:
									wrong_answer += 1
								answer_dict = {'id':choices.pk,'value':choices.value,'is_answer':choices.is_answer}
								chosen.append(answer_dict)
							row['answers'] = chosen
						else:
							row['answers'] = answers_chosen.text_answer

						if answers_chosen.document_image:
							row['image'] = str(answers_chosen.document_image)
						else:
							row['image'] = None
					if data['type'] == 'pdf':
						if wrong_answer > 0:
							for finding in findings:
								row['dupes'] = False
								if finding.value not in all_findings:
									row['dupes'] = True
									finding_question = finding.get_dict()
									possible_finding.append(finding_question)
									row['findings'] = possible_finding
									
									effects = Assessment_effect.objects.filter(question=question.pk,is_active=True)
									possible_effect = []
									for effect in effects:
										effect_question = effect.get_dict()
										possible_effect.append(effect_question)
									row['effects'] = possible_effect
									
									choices = Choice.objects.filter(question=question.pk,is_active=True)
									answers = []
									for choice in choices:
										answer_choice = choice.get_dict()
										answers.append(answer_choice)
									row['choices'] = answers

									image_answers = Assessment_upload_answer.objects.filter(is_deleted=False,question=question.pk,transaction_type=transaction_type.pk,company_assessment=data['id'])
									imageAnswersArr = []
									for image_answer in image_answers:
										imageAnswersArr.append(image_answer.get_dict())
									row['image_answers'] = imageAnswersArr
						datus.append(row)

						for all_finding in findings:
							if wrong_answer > 0:
								all_findings.append(all_finding.value)
					elif data['type'] == 'ppt':
						for finding in findings:
							# if finding.value not in all_findings:
							finding_question = finding.get_dict()
							possible_finding.append(finding_question)
							row['findings'] = possible_finding
								
						effects = Assessment_effect.objects.filter(question=question.pk,is_active=True)
						possible_effect = []
						for effect in effects:
							effect_question = effect.get_dict()
							possible_effect.append(effect_question)
						row['effects'] = possible_effect
						
						choices = Choice.objects.filter(question=question.pk,is_active=True)
						answers = []
						for choice in choices:
							answer_choice = choice.get_dict()
							answers.append(answer_choice)
						row['choices'] = answers

						image_answers = Assessment_upload_answer.objects.filter(is_deleted=False,question=question.pk,transaction_type=transaction_type.pk,company_assessment=data['id'])
						imageAnswersArr = []
						for image_answer in image_answers:
							image_data = image_answer.get_dict()
							imageAnswersArr.append(image_data)
						
						checkAnswerImage = Assessment_answer_image.objects.filter(question=question.pk,transaction_type=transaction_type.pk,company_assessment=data['id'])

						answerImageArr = []
						for img in checkAnswerImage:
							img_data = {}
							img_data['item_no'] = img.item_no
							img_data['answer_image'] = str(img.get_image())
							answerImageArr.append(img_data)
							
						row['answer_images'] = answerImageArr
						row['image_answers'] = imageAnswersArr
						datus.append(row)

						for all_finding in findings:
							if wrong_answer > 0:
								all_findings.append(all_finding.value)

			transactionTypesArr.append(transactionTypesRow)	

		results['transaction_types'] = transactionTypesArr
		results['scores'] = scoresArr
		results['data'] = datus
		return success_list(results,False)
	except Exception as e:
		exc_type, exc_obj, exc_tb = sys.exc_info()
		filename = os.path.split(exc_tb.tb_frame.f_code.co_filename)[1]
		linenum = sys.exc_traceback.tb_lineno,
		print(filename)
		print(linenum)
		print(e)
		return HttpResponse(e, status=400)

def generate(request):
	try:
		data = req_data(request)

		company_assessment = Company_assessment.objects.get(id=data['datus']['id'])

		prs = Presentation()

		title_slide_layout = prs.slide_layouts[0]
		slide = prs.slides.add_slide(title_slide_layout)
		title = slide.shapes.title
		subtitle = slide.placeholders[1]

		title.text = company_assessment.company.name
		subtitle.text = "Business Assessment"

		transaction_types = Exercise.objects.filter(id__in=data['datus']['transaction_type'])
		for transaction_type in transaction_types:
			transaction_type_title_slide_layout = prs.slide_layouts[0]
			transaction_type_slide = prs.slides.add_slide(transaction_type_title_slide_layout)
			transaction_type_title = transaction_type_slide.shapes.title
			transaction_type_subtitle = transaction_type_slide.placeholders[1]

			transaction_type_title.text = transaction_type.name
			transaction_type_subtitle.text = ""
			questions = Assessment_question.objects.filter(is_active=True,transaction_type=transaction_type.pk).order_by('transaction_type__name')
			# bullet_slide_layout1 = prs.slide_layouts[1]
			bullet_slide_layout = prs.slide_layouts[6]
			digit = Decimal(0)
			all_findings = []
			for question in questions:
				# slide3 = prs.slides.add_slide(bullet_slide_layout1)
				effects = Assessment_effect.objects.filter(is_active=True,question=question.pk)
				findings = Assessment_finding.objects.filter(is_active=True,question=question.pk)
				choices = Choice.objects.filter(is_active=True,question=question.pk)

				if not question.uploaded_question:
					try:
						answers_chosen = Assessment_answer.objects.get(company_assessment=data['datus']['id'],question=question.pk,is_deleted=False)
					except Assessment_answer.DoesNotExist:
						continue
					wrong_answer = Decimal(0)

					if answers_chosen.choice:
						for choices in answers_chosen.choice:
							choices = Choice.objects.get(pk=choices)
							if choices.is_answer == False:
								wrong_answer += 1


					if wrong_answer > 0:
						for finding in findings:
							if finding.value not in all_findings:
								slide2 = prs.slides.add_slide(bullet_slide_layout)

								### HEADER
								top2 = Inches(0)
								width2 = Inches(10)
								left2 = height2 = Inches(0)
								txBox2 = slide2.shapes.add_textbox(left2, top2, width2, height2)
								tf2 = txBox2.text_frame
								tf2.word_wrap = True

								q = tf2.add_paragraph()
								q.alignment = PP_ALIGN.CENTER
								run2 = q.add_run()
								run2.text = question.transaction_type.name
								font2 = run2.font
								font2.size = Pt(30)
								font2.bold = True
								####

								top = Inches(1)
								width = Inches(8)
								left = height = Inches(1)
								txBox = slide2.shapes.add_textbox(left, top, width, height)
								tf = txBox.text_frame
								tf.word_wrap = True

								o = tf.add_paragraph()
								o.alignment = PP_ALIGN.JUSTIFY
								run = o.add_run()
								digit += 1
								run.text = str(digit) + ". FINDING: " + finding.value
								font = run.font
								font.name = 'Calibri'
								font.size = Pt(18)
								font.bold = True
								font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

								for effect in effects:
									p = tf.add_paragraph()
									p.alignment = PP_ALIGN.JUSTIFY
									p.text = "POSSIBLE EFFECT: " + effect.value

								if answers_chosen.document_image:
									# img_slide = prs.slides.add_slide(bullet_slide_layout)
									img_path = 'web_admin'+settings.STATIC_URL+'uploads/'+str(answers_chosen.document_image)
									picture_top = Inches(4)
									picture_left = Inches(6)

									pic = slide2.shapes.add_picture(img_path, picture_left, picture_top)

						for all_finding in findings:
							all_findings.append(all_finding.value)
				else:
					try:
						answers_chosen = Assessment_answer.objects.get(uploaded_question=True,company_assessment=data['datus']['id'],question=question.pk,is_deleted=False)
					except Assessment_answer.DoesNotExist:
						pass
					wrong_answer = Decimal(0)

					# if answers_chosen.choice:
					# 	for choices in answers_chosen.choice:
					# 		choices = Choice.objects.get(pk=choices)
					# 		if choices.is_answer == False:
					# 			wrong_answer += 1

					image_answers = Assessment_upload_answer.objects.filter(is_deleted=False,question=question.pk,transaction_type=transaction_type.pk,company_assessment=data['datus']['id'])
					imageAnswersArr = []
					row = question.get_dict()

					for image_answer in image_answers:
						for correct_answers in row['answers']:
							if image_answer.item_no == correct_answers['item_no']:
								if image_answer.answer != correct_answers['answer']:
									wrong_answer += 1
						# imageAnswersArr.append(image_answer.get_dict())

					if wrong_answer > 0:
						for finding in findings:
							if finding.value not in all_findings:
								slide2 = prs.slides.add_slide(bullet_slide_layout)

								### HEADER
								top2 = Inches(0)
								width2 = Inches(10)
								left2 = height2 = Inches(0)
								txBox2 = slide2.shapes.add_textbox(left2, top2, width2, height2)
								tf2 = txBox2.text_frame
								tf2.word_wrap = True

								q = tf2.add_paragraph()
								q.alignment = PP_ALIGN.CENTER
								run2 = q.add_run()
								run2.text = question.transaction_type.name
								font2 = run2.font
								font2.size = Pt(30)
								font2.bold = True
								####

								top = Inches(1)
								width = Inches(8)
								left = height = Inches(1)
								txBox = slide2.shapes.add_textbox(left, top, width, height)
								tf = txBox.text_frame
								tf.word_wrap = True

								o = tf.add_paragraph()
								o.alignment = PP_ALIGN.JUSTIFY
								run = o.add_run()
								digit += 1
								run.text = str(digit) + ". FINDING: " + finding.value
								font = run.font
								font.name = 'Calibri'
								font.size = Pt(18)
								font.bold = True
								font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

								for effect in effects:
									p = tf.add_paragraph()
									p.alignment = PP_ALIGN.JUSTIFY
									p.text = "POSSIBLE EFFECT: " + effect.value

								# if answers_chosen.document_image:
								# 	img_path = 'web_admin'+settings.STATIC_URL+'uploads/'+str(answers_chosen.document_image)
								# 	picture_top = Inches(4)
								# 	picture_left = Inches(6)

								# 	pic = slide2.shapes.add_picture(img_path, picture_left, picture_top)

						for all_finding in findings:
							all_findings.append(all_finding.value)

		##################### SAVE CHOSEN RECOMMENDATION
		chosen_recommendations = []
		answer_recommendation = {}
		answer_recommendation['company_assessment'] = company_assessment.pk
		for recommendation in data['recommendations']:
			chosen_recommendations.append(recommendation['id'])

			slide_layout = prs.slide_layouts[1]

			slide = prs.slides.add_slide(slide_layout)
			title = slide.shapes.title
			subtitle = slide.placeholders[1]

			title.text = "Recommendations"
			subtitle.text = recommendation['value']

		answer_recommendation['recommendations'] = list_to_string(chosen_recommendations)
		try:
			instance = Generated_assessment_recommendation.objects.get(company_assessment=company_assessment.pk)
			generated_recommendation = Generated_assessment_recommendation_form(answer_recommendation,instance=instance)
		except Generated_assessment_recommendation.DoesNotExist:
			generated_recommendation = Generated_assessment_recommendation_form(answer_recommendation)

		if generated_recommendation.is_valid():
			generated_recommendation.save()
		#####################

		new_folder = os.path.join('web_admin'+settings.STATIC_URL, 'reports')

		try:
			os.mkdir(new_folder)
		except OSError as e:
			if e.errno != errno.EEXIST:
				pass
			else:
				print(e)

		company_assessment.is_generated = True
		company_assessment.save()

		prs.save('web_admin'+settings.STATIC_URL+'reports/'+company_assessment.reference_no+' - '+company_assessment.company.name+'.pptx')
		download_link = {'link' : str(settings.STATIC_URL)+'reports/'+company_assessment.reference_no+' - '+company_assessment.company.name+'.pptx'}
		return success_list(download_link,False)
	except Exception as e:
		return HttpResponse(e,status=400)

def generate_pdf(request,generate_report_id):
	company_assessment = Company_assessment.objects.get(id=generate_report_id)
	pagename = "Reference No.: " + company_assessment.reference_no + " - " + company_assessment.company.name
	datus = {
		"id" : company_assessment.pk,
		"reference_no" : str(company_assessment.reference_no),
		"company" : str(company_assessment.company.name),
		"transaction_type" : company_assessment.transaction_type,
		"status" : company_assessment.is_complete,
	}
	return render(request, 'generate_report/prints/generate_pdf.html', {"pagename":pagename,"datus":datus})

def read_chosen_recommendations(request):
	try:
		data = req_data(request)
		try:
			chosen_recommendations = Generated_assessment_recommendation.objects.get(company_assessment=data.get('id',None))
			recommendations = Assessment_recommendation.objects.filter(id__in=chosen_recommendations.recommendations,is_active=True)
		except Generated_assessment_recommendation.DoesNotExist:
			recommendations = {}
		results = {'data':[]}
		data = []
		for recommendation in recommendations:
			row = {}
			row['id'] = recommendation.pk
			row['value'] = recommendation.value
			data.append(row)

		results['data'] = data
		return success_list(results,False)
	except Exception as e:
		return HttpResponse(e,status=400)

def delete_report(request):
	try:
		data = req_data(request)

		if os.path.exists('web_admin'+data['link']):
			os.remove('web_admin'+data['link'])
		else:
			print("FILE DOES NOT EXISTS!!!!!!!!!!!!!!!!!!!!!!")

		return success()
	except Exception as e:
		return HttpResponse(e,status=400)

def new_score(request):
	try:
		data = req_data(request)

		try:
			checkAssessmentScore = Assessment_score.objects.get(company_assessment=data['company_assessment'],transaction_type=data['transaction_type'],question=data['question'],is_active=True)
		except Assessment_score.DoesNotExist:
			return error("This question has no score yet. Please be advised!")

		checkAssessmentScore.score = data['score']
		checkAssessmentScore.save()

		return success("Successfully saved!")
	except Exception as e:
		return HttpResponse(e,status=400)